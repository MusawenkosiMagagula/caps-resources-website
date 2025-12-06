import os
import shutil
import re
from pathlib import Path
import json

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from docx import Document
except ImportError:
    Document = None

# Configuration
RESOURCES_FOLDER = r"C:\caps-resources-website\Website"
ORGANIZED_FOLDER = r"C:\caps-resources-website\server\storage\pdfs"
THUMBNAILS_FOLDER = r"C:\caps-resources-website\images\products"

# Supported file types
SUPPORTED_EXTENSIONS = {
    'pdf': '.pdf',
    'word': ['.docx', '.doc'],
    'excel': ['.xlsx', '.xls'],
    'powerpoint': ['.pptx', '.ppt'],
    'archive': ['.zip', '.rar', '.7z']
}

# Grade mappings
GRADE_PATTERNS = {
    'preschool': ['preschool', 'pre-school', 'pre school', 'playgroup', 'creche'],
    'reception': ['reception', 'grade r', 'grade 0'],
    'grade1': ['grade 1', 'grade one', 'gr 1', 'gr1'],
    'grade2': ['grade 2', 'grade two', 'gr 2', 'gr2'],
    'grade3': ['grade 3', 'grade three', 'gr 3', 'gr3'],
    'grade4': ['grade 4', 'grade four', 'gr 4', 'gr4'],
    'grade5': ['grade 5', 'grade five', 'gr 5', 'gr5'],
    'grade6': ['grade 6', 'grade six', 'gr 6', 'gr6'],
    'grade7': ['grade 7', 'grade seven', 'gr 7', 'gr7'],
    'grade8': ['grade 8', 'grade eight', 'gr 8', 'gr8'],
    'grade9': ['grade 9', 'grade nine', 'gr 9', 'gr9'],
    'grade10': ['grade 10', 'grade ten', 'gr 10', 'gr10'],
    'grade11': ['grade 11', 'grade eleven', 'gr 11', 'gr11'],
    'grade12': ['grade 12', 'grade twelve', 'gr 12', 'gr12', 'matric']
}

# Subject mappings
SUBJECT_PATTERNS = {
    'Mathematics': ['mathematics', 'math', 'maths', 'wiskunde', 'numeracy'],
    'English': ['english', 'engels', 'home language english', 'hl english', 'literature', 'grammar'],
    'Afrikaans': ['afrikaans', 'afr', 'home language afrikaans', 'afr hl'],
    'Life Skills': ['life skills', 'lewensvaardighede', 'life orientation'],
    'Natural Sciences': ['natural sciences', 'science', 'natuurwetenskappe', 'natsci'],
    'Social Sciences': ['social sciences', 'history', 'geography', 'sosiale wetenskappe'],
    'Physical Sciences': ['physical sciences', 'fisiese wetenskappe', 'physics', 'chemistry'],
    'Life Sciences': ['life sciences', 'biology', 'lewenswetenskappe', 'biologie'],
    'Accounting': ['accounting', 'rekeningkunde'],
    'Business Studies': ['business studies', 'besigheidstudies'],
    'Economics': ['economics', 'ekonomie'],
    'Technology': ['technology', 'tegnologie', 'ict'],
    'Creative Arts': ['creative arts', 'arts', 'kreatiewe kunste', 'music', 'visual arts'],
    'Mathematical Literacy': ['mathematical literacy', 'math lit', 'maths literacy', 'wiskundige geletterdheid'],
    'First Additional Language': ['first additional language', 'fal', 'tweede taal'],
    'Home Language': ['home language', 'hl']
}

# Type mappings
TYPE_PATTERNS = {
    'worksheets': ['worksheet', 'work sheet', 'activity sheet', 'werksblad', 'exercise'],
    'assessments': ['assessment', 'test', 'exam', 'examination', 'toets', 'eksamen', 'quiz'],
    'lesson-plans': ['lesson plan', 'teaching plan', 'lesplan', 'lesson'],
    'activities': ['activity', 'activities', 'aktiwiteit', 'aktiwiteite', 'task'],
    'study-guides': ['study guide', 'revision', 'notes', 'studiegids', 'hersiening', 'summary']
}

def extract_text_from_pdf(pdf_path, max_pages=2):
    """Extract text from first few pages of PDF"""
    if PdfReader is None:
        return ""
    try:
        reader = PdfReader(pdf_path)
        text = ""
        for i in range(min(max_pages, len(reader.pages))):
            text += reader.pages[i].extract_text() + "\n"
        return text.lower()
    except Exception as e:
        print(f"    ⚠️  Error reading PDF text: {e}")
        return ""

def extract_text_from_docx(docx_path, max_paragraphs=20):
    """Extract text from Word document"""
    if Document is None:
        return ""
    try:
        doc = Document(docx_path)
        text = ""
        for i, para in enumerate(doc.paragraphs[:max_paragraphs]):
            if para.text.strip():
                text += para.text + "\n"
            if i >= max_paragraphs:
                break
        return text.lower()
    except Exception as e:
        print(f"    ⚠️  Error reading DOCX text: {e}")
        return ""

def extract_text_from_excel(excel_path, max_cells=100):
    """Extract text from Excel file"""
    if load_workbook is None:
        return ""
    try:
        wb = load_workbook(excel_path)
        ws = wb.active
        text = ""
        count = 0
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell and count < max_cells:
                    text += str(cell) + " "
                    count += 1
            if count >= max_cells:
                break
        return text.lower()
    except Exception as e:
        print(f"    ⚠️  Error reading Excel text: {e}")
        return ""

def extract_text_from_pptx(pptx_path, max_slides=3):
    """Extract text from PowerPoint"""
    if Presentation is None:
        return ""
    try:
        prs = Presentation(pptx_path)
        text = ""
        for i, slide in enumerate(prs.slides[:max_slides]):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.lower()
    except Exception as e:
        print(f"    ⚠️  Error reading PPTX text: {e}")
        return ""

def extract_text_from_file(file_path):
    """Extract text based on file type"""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.docx', '.doc']:
        return extract_text_from_docx(file_path)
    elif ext in ['.xlsx', '.xls']:
        return extract_text_from_excel(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    else:
        return ""

def get_file_type(file_path):
    """Determine the file type category"""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        return 'pdf'
    elif ext in ['.docx', '.doc']:
        return 'word'
    elif ext in ['.xlsx', '.xls']:
        return 'excel'
    elif ext in ['.pptx', '.ppt']:
        return 'powerpoint'
    elif ext in ['.zip', '.rar', '.7z']:
        return 'archive'
    else:
        return 'other'

def extract_grade(text, filename):
    """Extract grade from text and filename"""
    combined = text + " " + filename.lower()
    for grade, patterns in GRADE_PATTERNS.items():
        for pattern in patterns:
            if pattern.lower() in combined:
                return grade
    return None

def extract_subject(text, filename):
    """Extract subject from text and filename"""
    combined = text + " " + filename.lower()
    
    # Score each subject based on pattern matches
    subject_scores = {}
    for subject, patterns in SUBJECT_PATTERNS.items():
        score = 0
        for pattern in patterns:
            # Count how many times the pattern appears
            score += combined.count(pattern.lower())
        if score > 0:
            subject_scores[subject] = score
    
    # Return subject with highest score
    if subject_scores:
        return max(subject_scores, key=subject_scores.get)
    return "General"

def extract_type(text, filename):
    """Extract resource type from text and filename"""
    combined = text + " " + filename.lower()
    
    # Score each type
    type_scores = {}
    for resource_type, patterns in TYPE_PATTERNS.items():
        score = 0
        for pattern in patterns:
            score += combined.count(pattern.lower())
        if score > 0:
            type_scores[resource_type] = score
    
    if type_scores:
        return max(type_scores, key=type_scores.get)
    return "worksheets"

def extract_year(text, filename):
    """Extract year from text and filename"""
    combined = text + " " + filename
    year_pattern = r'\b(20\d{2})\b'
    matches = re.findall(year_pattern, combined)
    if matches:
        return matches[0]
    return "2024"

def clean_text_for_filename(text):
    """Clean text to be filename-safe"""
    # Remove special characters
    text = re.sub(r'[^\w\s-]', '', text)
    # Replace spaces with hyphens
    text = re.sub(r'\s+', '-', text)
    # Remove multiple hyphens
    text = re.sub(r'-+', '-', text)
    return text.strip('-').lower()

def get_file_size_readable(file_path):
    """Get file size in readable format"""
    size_bytes = os.path.getsize(file_path)
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024.0:
            return f"{size_bytes:.1f} {unit}"
        size_bytes /= 1024.0
    return f"{size_bytes:.1f} TB"

def get_page_count(file_path):
    """Get page/sheet count based on file type"""
    file_type = get_file_type(file_path)
    try:
        if file_type == 'pdf' and PdfReader:
            reader = PdfReader(file_path)
            return len(reader.pages)
        elif file_type == 'word' and Document:
            doc = Document(file_path)
            return len(doc.paragraphs)
        elif file_type == 'excel' and load_workbook:
            wb = load_workbook(file_path)
            return len(wb.sheetnames)
        elif file_type == 'powerpoint' and Presentation:
            prs = Presentation(file_path)
            return len(prs.slides)
    except:
        pass
    return 1

def process_file(file_path, output_base_dir):
    """Process a single file"""
    try:
        filename = os.path.basename(file_path)
        print(f"\nProcessing: {filename}")
        
        # Extract text from file
        text = extract_text_from_file(file_path)
        file_type = get_file_type(file_path)
        
        # Extract metadata
        grade = extract_grade(text, filename)
        subject = extract_subject(text, filename)
        resource_type = extract_type(text, filename)
        year = extract_year(text, filename)
        
        if not grade:
            print(f"  ⚠️  Could not determine grade, skipping...")
            return None
        
        # Get file info
        pages = get_page_count(file_path)
        file_size = get_file_size_readable(file_path)
        ext = os.path.splitext(filename)[1]
        
        # Create new filename
        new_filename = f"{grade}-{clean_text_for_filename(subject)}-{resource_type}-{year}{ext}"
        
        # Create output directory
        output_dir = os.path.join(output_base_dir, grade, subject.replace(' ', '-'))
        os.makedirs(output_dir, exist_ok=True)
        
        # Full output path
        output_path = os.path.join(output_dir, new_filename)
        
        # Handle duplicates
        counter = 1
        base_name = new_filename[:-len(ext)]
        while os.path.exists(output_path):
            new_filename = f"{base_name}-{counter}{ext}"
            output_path = os.path.join(output_dir, new_filename)
            counter += 1
        
        # Copy file
        shutil.copy2(file_path, output_path)
        
        print(f"  ✅ Organized as: {new_filename}")
        print(f"  📁 Type: {file_type.upper()} | Location: {output_dir}")
        print(f"  📊 Grade: {grade} | Subject: {subject} | Type: {resource_type}")
        print(f"  📄 Size: {file_size} | Pages: {pages}")
        
        return {
            'original_path': file_path,
            'new_path': output_path,
            'new_filename': new_filename,
            'file_type': file_type,
            'grade': grade,
            'subject': subject,
            'type': resource_type,
            'year': year,
            'pages': pages,
            'file_size': file_size,
            'extension': ext
        }
        
    except Exception as e:
        print(f"  ❌ Error processing {filename}: {e}")
        return None

def scan_and_organize():
    """Main function to scan and organize all files"""
    print("=" * 80)
    print("CAPS RESOURCES DOCUMENT ORGANIZER")
    print("=" * 80)
    
    # Create output directories
    os.makedirs(ORGANIZED_FOLDER, exist_ok=True)
    os.makedirs(THUMBNAILS_FOLDER, exist_ok=True)
    
    # Find all supported files
    supported_extensions = ['.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt']
    files_to_process = []
    
    if os.path.exists(RESOURCES_FOLDER):
        for root, dirs, files in os.walk(RESOURCES_FOLDER):
            for file in files:
                if any(file.lower().endswith(ext) for ext in supported_extensions):
                    files_to_process.append(os.path.join(root, file))
    else:
        print(f"❌ Resources folder not found: {RESOURCES_FOLDER}")
        return
    
    print(f"\n📂 Found {len(files_to_process)} supported files")
    print(f"📤 Output directory: {ORGANIZED_FOLDER}\n")
    
    # Process each file
    results = []
    successful = 0
    failed = 0
    
    for file_path in files_to_process:
        result = process_file(file_path, ORGANIZED_FOLDER)
        if result:
            results.append(result)
            successful += 1
        else:
            failed += 1
    
    # Summary
    print("\n" + "=" * 80)
    print("SUMMARY")
    print("=" * 80)
    print(f"✅ Successfully organized: {successful}")
    print(f"❌ Failed: {failed}")
    print(f"📊 Total processed: {len(files_to_process)}")
    
    # Generate summary by grade and subject
    print("\n📚 Resources by Grade:")
    by_grade = {}
    for result in results:
        grade = result['grade']
        if grade not in by_grade:
            by_grade[grade] = []
        by_grade[grade].append(result)
    
    for grade in sorted(by_grade.keys()):
        print(f"\n  {grade.upper()}: {len(by_grade[grade])} resources")
        subjects = {}
        for item in by_grade[grade]:
            subject = item['subject']
            subjects[subject] = subjects.get(subject, 0) + 1
        for subject, count in sorted(subjects.items()):
            print(f"    - {subject}: {count}")
    
    # File type breakdown
    print("\n📋 Resources by File Type:")
    by_type = {}
    for result in results:
        ftype = result['file_type'].upper()
        by_type[ftype] = by_type.get(ftype, 0) + 1
    for ftype, count in sorted(by_type.items()):
        print(f"  - {ftype}: {count}")
    
    # Save results
    results_file = os.path.join(ORGANIZED_FOLDER, '_organization_results.json')
    with open(results_file, 'w', encoding='utf-8') as f:
        json.dump(results, f, indent=2, ensure_ascii=False)
    
    print(f"\n💾 Results saved to: {results_file}")
    print("\n✨ Organization complete!")

if __name__ == "__main__":
    scan_and_organize()
